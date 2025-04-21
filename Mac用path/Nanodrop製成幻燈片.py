from pptx import Presentation
from pptx.util import Inches, Cm
import os

# 設定圖片所在的資料夾和輸出 PowerPoint 的名稱
image_folder = "/Users/jasontingting/Desktop/Nanodrop/2024-12-14 K1-K4 NGS PCR 精製/折線圖保存"  # 輸出圖片的資料夾
output_pptx = "/Users/jasontingting/Desktop/Nanodrop/2024-12-14 K1-K4 NGS PCR 精製/折線圖保存/output_presentation.pptx"  # PowerPoint 檔案名稱

# 初始化 PowerPoint 文件
prs = Presentation()

# 讀取資料夾中的所有圖片
image_files = [f for f in os.listdir(image_folder) if f.endswith(('.png', '.jpg', '.jpeg'))]

# 設定每頁最多顯示 20 張圖片
images_per_slide = 20
current_count = 0

# 設定每張圖片的位置和大小
pic_width = Cm(5)  # 設定圖片寬度為 5 厘米
pic_height = Cm(3)  # 設定圖片高度為 3 厘米
horizontal_gap = Inches(0.2)  # 圖片之間的水平間隙
vertical_gap = Inches(0.2)  # 圖片之間的垂直間隙
max_cols = 5  # 每行最多放置的圖片數量
max_rows = 4  # 每頁最多放置的行數

# 幻燈片的寬度和高度
slide_width = prs.slide_width  # 幻燈片的寬度 (以像素為單位)
slide_height = prs.slide_height  # 幻燈片的高度 (以像素為單位)

# 計算居中位置
total_width = pic_width * max_cols + horizontal_gap * (max_cols - 1)  # 所有圖片的總寬度
total_height = pic_height * max_rows + vertical_gap * (max_rows - 1)  # 所有圖片的總高度

# 計算圖片左上角的起始位置
start_x = (slide_width - total_width) / 2  # 水平居中
start_y = (slide_height - total_height) / 2  # 垂直居中

# 逐張圖片插入 PowerPoint 幻燈片
for idx, image in enumerate(image_files):
    if current_count % images_per_slide == 0:
        # 每20張圖片創建新幻燈片
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白幻燈片模板
        current_x = start_x  # 初始 X 位置（居中）
        current_y = start_y  # 初始 Y 位置（居中）

    # 插入圖片
    img_path = os.path.join(image_folder, image)
    pic = slide.shapes.add_picture(img_path, current_x, current_y, width=pic_width, height=pic_height)
    
    # 更新X和Y位置以便下一張圖片
    current_x += pic_width + horizontal_gap
    
    # 如果一行放不下，換到下一行
    if (current_x + pic_width) > slide_width:  # 如果圖片超過了幻燈片的寬度
        current_x = start_x  # 重置為居中位置
        current_y += pic_height + vertical_gap  # 移到下一行
    
    current_count += 1

# 保存簡報到指定路徑
prs.save(output_pptx)

print(f"所有圖表已保存至資料夾: {output_pptx}") 